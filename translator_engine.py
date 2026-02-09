from __future__ import annotations

from abc import ABC, abstractmethod
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from io import BytesIO
import re
from typing import Callable, Iterable, Optional

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


TranslateOne = Callable[[str], str]
ProgressCallback = Callable[[int, int], None]

MAX_BATCH_ITEMS = 12
MAX_BATCH_CHARS = 2000
ITEM_START = "<<<ITEM"
ITEM_END = "<<<END"


class TranslationStrategy(ABC):
    @abstractmethod
    def process(
        self,
        file_bytes: bytes,
        source_lang: str,
        target_lang: str,
        translate_one: TranslateOne,
        progress_cb: Optional[ProgressCallback] = None,
    ) -> BytesIO:
        raise NotImplementedError


@dataclass
class TranslatorEngine:
    strategies: dict[str, TranslationStrategy]

    def process(
        self,
        filename: str,
        file_bytes: bytes,
        source_lang: str,
        target_lang: str,
        translate_one: TranslateOne,
        progress_cb: Optional[ProgressCallback] = None,
    ) -> BytesIO:
        ext = _get_extension(filename)
        if ext not in self.strategies:
            raise ValueError(f"Unsupported file type: {ext}")
        return self.strategies[ext].process(
            file_bytes=file_bytes,
            source_lang=source_lang,
            target_lang=target_lang,
            translate_one=translate_one,
            progress_cb=progress_cb,
        )


class DocxStrategy(TranslationStrategy):
    def process(
        self,
        file_bytes: bytes,
        source_lang: str,
        target_lang: str,
        translate_one: TranslateOne,
        progress_cb: Optional[ProgressCallback] = None,
    ) -> BytesIO:
        doc = Document(BytesIO(file_bytes))
        paragraphs = _collect_docx_paragraphs(doc)
        _translate_paragraphs_in_place(
            paragraphs=paragraphs,
            translate_one=translate_one,
            progress_cb=progress_cb,
        )
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output


class PptxStrategy(TranslationStrategy):
    def process(
        self,
        file_bytes: bytes,
        source_lang: str,
        target_lang: str,
        translate_one: TranslateOne,
        progress_cb: Optional[ProgressCallback] = None,
    ) -> BytesIO:
        prs = Presentation(BytesIO(file_bytes))
        items = _collect_pptx_paragraphs(prs)
        _translate_pptx_items_in_place(
            items=items,
            translate_one=translate_one,
            progress_cb=progress_cb,
        )
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output


@dataclass
class PptxParagraphItem:
    paragraph: object
    original_text: str
    base_font_size: float
    slide_index: int


def _get_extension(filename: str) -> str:
    lower = filename.lower()
    if "." not in lower:
        return ""
    return lower.rsplit(".", 1)[-1]


def _collect_docx_paragraphs(doc: Document) -> list:
    paragraphs = [p for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph)
    return paragraphs


def _iter_shapes(shapes: Iterable) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _get_base_font_size(text_frame) -> float:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                return run.font.size.pt
    return 18.0


def _collect_pptx_paragraphs(prs: Presentation) -> list[PptxParagraphItem]:
    items: list[PptxParagraphItem] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            base_size = _get_base_font_size(shape.text_frame)
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                items.append(
                    PptxParagraphItem(
                        paragraph=paragraph,
                        original_text=text,
                        base_font_size=base_size,
                        slide_index=slide_index,
                    )
                )
    return items


def _translate_paragraphs_in_place(
    paragraphs: list,
    translate_one: TranslateOne,
    progress_cb: Optional[ProgressCallback],
) -> None:
    texts = [p.text for p in paragraphs]
    translated = _translate_texts_batched(texts, translate_one)
    total = len(paragraphs)
    for idx, (paragraph, new_text) in enumerate(zip(paragraphs, translated), start=1):
        if new_text:
            paragraph.text = new_text
        if progress_cb:
            progress_cb(idx, total)


def _translate_pptx_items_in_place(
    items: list[PptxParagraphItem],
    translate_one: TranslateOne,
    progress_cb: Optional[ProgressCallback],
) -> None:
    texts = [item.original_text for item in items]
    translated = _translate_texts_batched(texts, translate_one)
    total = len(items)
    for idx, (item, new_text) in enumerate(zip(items, translated), start=1):
        if new_text:
            item.paragraph.text = new_text
            _maybe_shrink_runs(
                item.paragraph,
                item.base_font_size,
                item.original_text,
                new_text,
            )
        if progress_cb:
            progress_cb(idx, total)


def _translate_texts_concurrently(
    texts: list[str],
    translate_one: TranslateOne,
) -> list[Optional[str]]:
    if not texts:
        return []

    def _safe_translate(text: str) -> Optional[str]:
        if not text.strip():
            return None
        try:
            return translate_one(text)
        except Exception as exc:
            print(f"翻译失败，跳过该段落: {exc}")
            return None

    with ThreadPoolExecutor() as executor:
        results = list(executor.map(_safe_translate, texts))
    return results


def _translate_texts_batched(
    texts: list[str],
    translate_one: TranslateOne,
) -> list[Optional[str]]:
    if not texts:
        return []

    batches: list[list[tuple[int, str]]] = []
    current: list[tuple[int, str]] = []
    current_chars = 0
    for idx, text in enumerate(texts):
        text_len = len(text)
        if current and (
            len(current) >= MAX_BATCH_ITEMS
            or current_chars + text_len > MAX_BATCH_CHARS
        ):
            batches.append(current)
            current = []
            current_chars = 0
        current.append((idx, text))
        current_chars += text_len
    if current:
        batches.append(current)

    results: list[Optional[str]] = [None] * len(texts)
    for batch in batches:
        batch_results = _translate_batch(batch, translate_one)
        for idx, translated in batch_results:
            results[idx] = translated
    return results


def _translate_batch(
    batch: list[tuple[int, str]],
    translate_one: TranslateOne,
) -> list[tuple[int, Optional[str]]]:
    prompt_lines = [
        "Translate each item separately. Keep the tags unchanged.",
        "Return only the translated items in the same tag format.",
    ]
    for i, (_, text) in enumerate(batch, start=1):
        prompt_lines.append(f"{ITEM_START} {i}>>>")
        prompt_lines.append(text)
        prompt_lines.append(f"{ITEM_END} {i}>>>")
    prompt = "\n".join(prompt_lines)

    try:
        response = translate_one(prompt)
    except Exception as exc:
        print(f"批处理翻译失败，回退为逐条翻译: {exc}")
        return _fallback_translate(batch, translate_one)

    parsed = _parse_batch_response(response, len(batch))
    if parsed is None:
        return _fallback_translate(batch, translate_one)

    results: list[tuple[int, Optional[str]]] = []
    for (idx, _), translated in zip(batch, parsed):
        results.append((idx, translated))
    return results


def _parse_batch_response(response: str, expected: int) -> Optional[list[str]]:
    pattern = re.compile(
        r"<<<ITEM\s+(\d+)>>>\s*(.*?)\s*<<<END\s+\1>>>",
        re.DOTALL,
    )
    matches = pattern.findall(response or "")
    if len(matches) != expected:
        return None
    results = [None] * expected
    for idx_str, content in matches:
        try:
            idx = int(idx_str)
        except ValueError:
            return None
        if idx < 1 or idx > expected:
            return None
        results[idx - 1] = content.strip()
    if any(item is None for item in results):
        return None
    return results  # type: ignore[return-value]


def _fallback_translate(
    batch: list[tuple[int, str]],
    translate_one: TranslateOne,
) -> list[tuple[int, Optional[str]]]:
    results: list[tuple[int, Optional[str]]] = []
    for idx, text in batch:
        try:
            results.append((idx, translate_one(text)))
        except Exception as exc:
            print(f"翻译失败，跳过该段落: {exc}")
            results.append((idx, None))
    return results


def _maybe_shrink_runs(paragraph, base_size: float, original: str, translated: str) -> None:
    if not original or not translated:
        return
    if len(translated) <= len(original) * 1.2:
        return
    for run in paragraph.runs:
        size = run.font.size.pt if run.font.size else base_size
        new_size = max(10.0, size * 0.9)
        run.font.size = Pt(new_size)
